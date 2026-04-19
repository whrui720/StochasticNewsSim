"""Build the project presentation .pptx (academic, professional).

Math formatting: strings may contain $...$ blocks with LaTeX-ish markup.
Supported macros:
  \\mu \\lambda \\alpha \\beta \\sigma \\rho \\pi \\tau \\theta \\phi
  \\psi \\chi \\epsilon \\delta \\gamma \\omega \\Gamma \\Delta \\Lambda
  \\Sigma \\Pi ...
  \\cdot \\times \\to \\approx \\leq \\geq \\neq \\pm \\infty \\ell
  \\sim \\propto \\equiv \\ll \\gg \\in \\subset \\subseteq
  \\sum \\prod \\int \\partial
  \\log \\exp \\max \\min \\arg \\argmax \\sin \\cos
  _x / _{abc}, ^x / ^{abc}
  \\binom{a}{b}   (rendered as \"C(a, b)\")
  \\frac{a}{b}    (rendered inline as a/b)
  \\mathbb{R}     (blackboard bold, R/N/Z/Q/C/E/P)
  \\mathrm{text}  (upright multi-letter)
  \\text{...}     (upright, same as mathrm)
  \\{  \\}        (literal braces)
  - (minus), * (multiplication dot)
Single letters in math render italic; multi-letter sequences upright.
"""
from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

ROOT = Path(__file__).parent
PARTA = ROOT / "partAresults"
PARTB = ROOT / "partBresults"
OUT = ROOT / "StochasticNewsSim_presentation.pptx"

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
SLATE = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
LIGHT_GREY = RGBColor(0x77, 0x77, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_SUB = RGBColor(0xCF, 0xD8, 0xE5)
BOX_BG = RGBColor(0xF4, 0xF6, 0xFA)

TEXT_FONT = "Calibri"
MATH_FONT = "Cambria Math"

# ---------------------------------------------------------------------------
# Math renderer
# ---------------------------------------------------------------------------
GREEK_LOWER = {
    "alpha": "\u03b1", "beta": "\u03b2", "gamma": "\u03b3", "delta": "\u03b4",
    "epsilon": "\u03b5", "varepsilon": "\u03b5", "zeta": "\u03b6",
    "eta": "\u03b7", "theta": "\u03b8", "iota": "\u03b9", "kappa": "\u03ba",
    "lambda": "\u03bb", "mu": "\u03bc", "nu": "\u03bd", "xi": "\u03be",
    "pi": "\u03c0", "rho": "\u03c1", "sigma": "\u03c3", "tau": "\u03c4",
    "upsilon": "\u03c5", "phi": "\u03c6", "chi": "\u03c7", "psi": "\u03c8",
    "omega": "\u03c9",
}
GREEK_UPPER = {
    "Gamma": "\u0393", "Delta": "\u0394", "Theta": "\u0398",
    "Lambda": "\u039b", "Xi": "\u039e", "Pi": "\u03a0",
    "Sigma": "\u03a3", "Phi": "\u03a6", "Psi": "\u03a8", "Omega": "\u03a9",
}
SYMBOLS = {
    "cdot": "\u00b7", "times": "\u00d7", "to": "\u2192",
    "approx": "\u2248", "leq": "\u2264", "geq": "\u2265",
    "neq": "\u2260", "pm": "\u00b1", "infty": "\u221e", "ell": "\u2113",
    "sim": "\u223c", "propto": "\u221d", "equiv": "\u2261",
    "ll": "\u226a", "gg": "\u226b", "in": "\u2208",
    "subset": "\u2282", "subseteq": "\u2286",
    "sum": "\u2211", "prod": "\u220f", "int": "\u222b", "partial": "\u2202",
    "forall": "\u2200", "exists": "\u2203", "mid": "|", "vert": "|",
    "uparrow": "\u2191", "downarrow": "\u2193",
}
FUNCS = {"log", "exp", "max", "min", "argmax", "argmin", "arg",
         "sin", "cos", "tan", "det", "var", "cov", "Var", "Cov", "E"}
BB = {"R": "\u211d", "N": "\u2115", "Z": "\u2124", "Q": "\u211a",
      "C": "\u2102", "E": "\U0001D53C", "P": "\u2119"}


def _add_run(p, text, *, size, color, italic=False, bold=False,
             font=TEXT_FONT, baseline=None):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = italic
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    if baseline is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("baseline", str(baseline))
    return run


def _find_matching_brace(s, i):
    depth = 1
    j = i + 1
    while j < len(s) and depth:
        if s[j] == "{":
            depth += 1
        elif s[j] == "}":
            depth -= 1
            if depth == 0:
                return j
        j += 1
    return len(s)


def _render_math_group(p, expr, *, size, color, baseline=None):
    i = 0
    n = len(expr)
    while i < n:
        c = expr[i]
        if c == "\\":
            # Escaped braces \{ and \}
            if i + 1 < n and expr[i + 1] in "{}":
                _add_run(p, expr[i + 1], size=size, color=color,
                         font=MATH_FONT, baseline=baseline)
                i += 2
                continue
            m = re.match(r"\\([a-zA-Z]+)", expr[i:])
            if m:
                name = m.group(1)
                consumed = 1 + len(name)
                # macros that eat further {...} args
                if name == "binom":
                    pos = i + consumed
                    if pos < n and expr[pos] == "{":
                        j1 = _find_matching_brace(expr, pos)
                        a = expr[pos + 1:j1]
                        pos2 = j1 + 1
                        if pos2 < n and expr[pos2] == "{":
                            j2 = _find_matching_brace(expr, pos2)
                            b = expr[pos2 + 1:j2]
                            _add_run(p, "C", size=size, color=color,
                                     italic=True, font=MATH_FONT,
                                     baseline=baseline)
                            _add_run(p, "(", size=size, color=color,
                                     font=MATH_FONT, baseline=baseline)
                            _render_math_group(p, a, size=size, color=color,
                                               baseline=baseline)
                            _add_run(p, ", ", size=size, color=color,
                                     font=MATH_FONT, baseline=baseline)
                            _render_math_group(p, b, size=size, color=color,
                                               baseline=baseline)
                            _add_run(p, ")", size=size, color=color,
                                     font=MATH_FONT, baseline=baseline)
                            i = j2 + 1
                            continue
                if name == "frac":
                    pos = i + consumed
                    if pos < n and expr[pos] == "{":
                        j1 = _find_matching_brace(expr, pos)
                        a = expr[pos + 1:j1]
                        pos2 = j1 + 1
                        if pos2 < n and expr[pos2] == "{":
                            j2 = _find_matching_brace(expr, pos2)
                            b = expr[pos2 + 1:j2]
                            _render_math_group(p, a, size=size, color=color,
                                               baseline=baseline)
                            _add_run(p, "\u2009/\u2009", size=size,
                                     color=color, font=MATH_FONT,
                                     baseline=baseline)
                            _render_math_group(p, b, size=size, color=color,
                                               baseline=baseline)
                            i = j2 + 1
                            continue
                if name in ("mathrm", "text", "operatorname"):
                    pos = i + consumed
                    if pos < n and expr[pos] == "{":
                        j1 = _find_matching_brace(expr, pos)
                        content = expr[pos + 1:j1]
                        _add_run(p, content, size=size, color=color,
                                 font=MATH_FONT, baseline=baseline)
                        i = j1 + 1
                        continue
                if name == "mathbb":
                    pos = i + consumed
                    if pos < n and expr[pos] == "{":
                        j1 = _find_matching_brace(expr, pos)
                        content = expr[pos + 1:j1]
                        out = "".join(BB.get(ch, ch) for ch in content)
                        _add_run(p, out, size=size, color=color,
                                 font=MATH_FONT, baseline=baseline)
                        i = j1 + 1
                        continue
                if name == "hat":
                    pos = i + consumed
                    if pos < n and expr[pos] == "{":
                        j1 = _find_matching_brace(expr, pos)
                        inner = expr[pos + 1:j1]
                        _render_math_group(p, inner, size=size, color=color,
                                           baseline=baseline)
                        _add_run(p, "\u0302", size=size, color=color,
                                 font=MATH_FONT, baseline=baseline)
                        i = j1 + 1
                        continue
                if name in GREEK_LOWER:
                    _add_run(p, GREEK_LOWER[name], size=size, color=color,
                             italic=True, font=MATH_FONT, baseline=baseline)
                elif name in GREEK_UPPER:
                    _add_run(p, GREEK_UPPER[name], size=size, color=color,
                             font=MATH_FONT, baseline=baseline)
                elif name in SYMBOLS:
                    _add_run(p, SYMBOLS[name], size=size, color=color,
                             font=MATH_FONT, baseline=baseline)
                elif name in FUNCS:
                    _add_run(p, name, size=size, color=color,
                             font=MATH_FONT, baseline=baseline)
                else:
                    _add_run(p, "\\" + name, size=size, color=color,
                             font=MATH_FONT, baseline=baseline)
                i += consumed
                continue
        if c == "_":
            i += 1
            if i < n and expr[i] == "{":
                j = _find_matching_brace(expr, i)
                content = expr[i + 1:j]
                _render_sub_sup(p, content, size=size, color=color,
                                is_sup=False, outer_baseline=baseline)
                i = j + 1
            elif i < n:
                _render_sub_sup(p, expr[i], size=size, color=color,
                                is_sup=False, outer_baseline=baseline)
                i += 1
            continue
        if c == "^":
            i += 1
            if i < n and expr[i] == "{":
                j = _find_matching_brace(expr, i)
                content = expr[i + 1:j]
                _render_sub_sup(p, content, size=size, color=color,
                                is_sup=True, outer_baseline=baseline)
                i = j + 1
            elif i < n:
                _render_sub_sup(p, expr[i], size=size, color=color,
                                is_sup=True, outer_baseline=baseline)
                i += 1
            continue
        if c.isalpha():
            j = i
            while j < n and expr[j].isalpha():
                j += 1
            word = expr[i:j]
            if len(word) == 1:
                _add_run(p, word, size=size, color=color, italic=True,
                         font=MATH_FONT, baseline=baseline)
            else:
                _add_run(p, word, size=size, color=color,
                         font=MATH_FONT, baseline=baseline)
            i = j
            continue
        if c.isdigit() or c == ".":
            j = i
            while j < n and (expr[j].isdigit() or expr[j] == "."):
                j += 1
            _add_run(p, expr[i:j], size=size, color=color,
                     font=MATH_FONT, baseline=baseline)
            i = j
            continue
        if c == "-":
            _add_run(p, "\u2212", size=size, color=color,
                     font=MATH_FONT, baseline=baseline)
            i += 1
            continue
        if c == "*":
            _add_run(p, "\u00b7", size=size, color=color,
                     font=MATH_FONT, baseline=baseline)
            i += 1
            continue
        _add_run(p, c, size=size, color=color,
                 font=MATH_FONT, baseline=baseline)
        i += 1


def _render_sub_sup(p, content, *, size, color, is_sup, outer_baseline=None):
    inner_size = size * 0.72
    shift = 30000 if is_sup else -25000
    if outer_baseline is not None:
        shift += outer_baseline
    is_abbrev = content.isalpha() and len(content) > 1
    if is_abbrev:
        _add_run(p, content, size=inner_size, color=color,
                 font=MATH_FONT, baseline=shift)
    else:
        _render_math_group(p, content, size=inner_size, color=color,
                           baseline=shift)


def _split_math(s):
    parts = []
    i = 0
    while i < len(s):
        if s[i] == "$":
            j = s.find("$", i + 1)
            if j == -1:
                parts.append(("t", s[i:]))
                break
            parts.append(("m", s[i + 1:j]))
            i = j + 1
        else:
            j = s.find("$", i)
            if j == -1:
                parts.append(("t", s[i:]))
                break
            parts.append(("t", s[i:j]))
            i = j
    return parts


def _render_string(p, s, *, size, color, bold=False, font=TEXT_FONT):
    for kind, content in _split_math(s):
        if kind == "m":
            _render_math_group(p, content, size=size, color=color)
        else:
            if not content:
                continue
            _add_run(p, content, size=size, color=color,
                     bold=bold, font=font)


# ---------------------------------------------------------------------------
# Shape / layout helpers
# ---------------------------------------------------------------------------
def add_rect_band(slide, top, height, color, left=0, width=None):
    if width is None:
        width = SW
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    return s


def add_plain_text(slide, left, top, width, height, text, *, size=18,
                   bold=False, color=SLATE, align=PP_ALIGN.LEFT,
                   font=TEXT_FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text if isinstance(text, list) else text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _render_string(p, line, size=size, color=color, bold=bold, font=font)
    return tb


def add_rich_bullets(slide, left, top, width, height, bullets, *, size=16,
                     color=SLATE, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(bullets):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        marker = "\u2022  " if level == 0 else "\u2013  "
        eff_size = size - level * 2
        _add_run(p, marker, size=eff_size, color=color, font=TEXT_FONT)
        _render_string(p, text, size=eff_size, color=color, font=TEXT_FONT)
    return tb


def add_rich_paragraph(slide, left, top, width, height, text, *, size=18,
                       bold=False, color=SLATE, align=PP_ALIGN.LEFT,
                       font=TEXT_FONT):
    return add_plain_text(slide, left, top, width, height, text,
                          size=size, bold=bold, color=color, align=align,
                          font=font)


def add_image_fit(slide, path, left, top, max_w, max_h):
    with Image.open(path) as im:
        w, h = im.size
    ratio = min(max_w / w, max_h / h)
    fw, fh = int(w * ratio), int(h * ratio)
    pic_l = left + (max_w - fw) // 2
    pic_t = top + (max_h - fh) // 2
    slide.shapes.add_picture(str(path), pic_l, pic_t, width=fw, height=fh)


def slide_header(slide, title, subtitle=None):
    add_rect_band(slide, 0, Inches(0.9), NAVY)
    add_plain_text(slide, Inches(0.4), Inches(0.16), Inches(12.5),
                   Inches(0.5), title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_plain_text(slide, Inches(0.4), Inches(0.56), Inches(12.5),
                       Inches(0.32), subtitle, size=13, color=HEADER_SUB)
    add_rect_band(slide, Inches(7.15), Inches(0.04), ACCENT)
    add_plain_text(slide, Inches(0.4), Inches(7.20), Inches(12.5),
                   Inches(0.25),
                   "Stochastics of News Article Propagation  \u2014  "
                   "Wang \u00b7 Madugula \u00b7 Thavamani",
                   size=10, color=LIGHT_GREY)


def make_table(slide, left, top, width, height, headers, rows,
               highlight_row=None, header_size=12, body_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                   height).table
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = ""
        p = c.text_frame.paragraphs[0]
        _render_string(p, h, size=header_size, color=WHITE, bold=True)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            p = c.text_frame.paragraphs[0]
            is_hl = (highlight_row is not None and i == highlight_row)
            col = ACCENT if is_hl else SLATE
            _render_string(p, val, size=body_size, color=col, bold=is_hl)
    return table


def add_formula_box(slide, left, top, width, height, title, body_items, *,
                    title_size=12, body_size=11, formula_size=13,
                    header_color=NAVY, body_bg=BOX_BG,
                    title_color=WHITE, text_color=SLATE, accent=ACCENT,
                    line_spacing=1.18):
    """`body_items` is a list where each item is either a string (bullet) or
    a 2-tuple (text, kind) with kind in {"b" (bullet), "f" (centered formula),
    "t" (centered text block), "sub" (small centered caption)}."""
    outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                   width, height)
    outer.line.color.rgb = header_color
    outer.line.width = Pt(0.75)
    outer.fill.solid()
    outer.fill.fore_color.rgb = body_bg
    outer.text_frame.text = ""
    hh = Inches(0.36)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, hh)
    hdr.line.fill.background()
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_color
    hdr.text_frame.text = ""
    tb_hdr = slide.shapes.add_textbox(left + Inches(0.12),
                                      top + Inches(0.04),
                                      width - Inches(0.24), hh)
    tb_hdr.text_frame.word_wrap = False
    tb_hdr.text_frame.margin_left = Inches(0.02)
    tb_hdr.text_frame.margin_top = Inches(0.0)
    p_hdr = tb_hdr.text_frame.paragraphs[0]
    _render_string(p_hdr, title, size=title_size, color=title_color,
                   bold=True)

    body_left = left + Inches(0.18)
    body_top = top + hh + Inches(0.08)
    body_w = width - Inches(0.36)
    body_h = height - hh - Inches(0.14)
    tb = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(body_items):
        if isinstance(item, tuple):
            text, kind = item
        else:
            text, kind = item, "b"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        if kind == "f":
            p.alignment = PP_ALIGN.CENTER
            _render_string(p, text, size=formula_size, color=NAVY)
        elif kind == "t":
            p.alignment = PP_ALIGN.CENTER
            _render_string(p, text, size=body_size, color=text_color)
        elif kind == "sub":
            p.alignment = PP_ALIGN.CENTER
            _render_string(p, text, size=body_size - 1, color=LIGHT_GREY)
        else:  # bullet
            p.alignment = PP_ALIGN.LEFT
            _add_run(p, "\u2022  ", size=body_size, color=text_color,
                     font=TEXT_FONT)
            _render_string(p, text, size=body_size, color=text_color)
    return outer


# ---------------------------------------------------------------------------
# Presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------------- Slide 1: Title ----------------
s = prs.slides.add_slide(BLANK)
add_rect_band(s, 0, SH, RGBColor(0xFA, 0xFA, 0xFC))
add_rect_band(s, Inches(2.9), Inches(0.06), NAVY)
add_rect_band(s, Inches(4.6), Inches(0.02), ACCENT)
add_plain_text(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5),
               "STAT 460  \u00b7  Stochastic Processes  \u00b7  Final Project",
               size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_plain_text(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.4),
               "Stochastics of News Article Propagation",
               size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_plain_text(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.5),
               "Modelling counts and reliability dynamics around the 2017 "
               "Las Vegas shooting",
               size=18, color=SLATE, align=PP_ALIGN.CENTER)
add_plain_text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
               "Harry Wang  \u00b7  Sriyan Madugula  \u00b7  Ruthesh Thavamani",
               size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_plain_text(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.4),
               "University of Michigan",
               size=14, color=SLATE, align=PP_ALIGN.CENTER)

# ---------------- Slide 2: Introduction ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Introduction",
             "Joint stochastic modelling of news volume and reliability "
             "after a singular event")

add_rich_paragraph(s, Inches(0.45), Inches(1.05), Inches(6.4), Inches(0.4),
                   "Problem", size=14, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(1.45), Inches(6.4), Inches(2.1), [
    "A single high-impact event triggers a cascade of articles across "
    "many outlets over the following days/weeks",
    "Coverage volume decays, but the *mix* of reliable vs unreliable "
    "sources also evolves \u2014 and is itself a stochastic process",
    "Test case: 2017 Las Vegas shooting \u2014 clean, isolated, "
    "exogenously-triggered news shock",
], size=12)

add_rich_paragraph(s, Inches(0.45), Inches(3.55), Inches(6.4), Inches(0.4),
                   "Questions and gap", size=14, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(3.95), Inches(6.4), Inches(3.1), [
    "Part A:  what stochastic process best describes the daily article "
    "count $N(t)$ following the event?",
    "Part B:  what process governs the reliable / unreliable composition "
    "$\\pi_{pos}(t)$ of that coverage?",
    "Gap addressed",
    ("Volume dynamics and reliability dynamics are typically studied in "
     "isolation \u2014 we model both on the same event window", 1),
    ("Stochastic-process treatment of the reliability *mix* at "
     "event-level granularity is rare", 1),
    ("Cross-excitation between reliable and unreliable streams is "
     "largely unexamined", 1),
], size=12)

# Right column: literature box
add_formula_box(s, Inches(7.0), Inches(1.05), Inches(6.1), Inches(6.05),
                "Prior work",
                [
                    ("Point-process models of online attention", "t"),
                    ("Crane & Sornette (PNAS, 2008) \u2014 endogenous vs "
                     "exogenous bursts in YouTube views via Hawkes-type "
                     "response functions", "b"),
                    ("Zhao et al. SEISMIC (KDD, 2015);  Mishra, Rizoiu & "
                     "Xie (2016) \u2014 self-exciting models for tweet and "
                     "retweet cascades", "b"),
                    ("Daley & Vere-Jones (2003) \u2014 standard reference "
                     "for inhomogeneous Poisson and Hawkes processes", "b"),
                    ("Misinformation and reliability", "t"),
                    ("Vosoughi, Roy & Aral (Science, 2018) \u2014 false "
                     "news spreads farther and faster than true news on "
                     "Twitter (descriptive, not generative)", "b"),
                    ("Allcott & Gentzkow (JEP, 2017) \u2014 economics of "
                     "fake news during the 2016 election", "b"),
                    ("Tambuscio et al. (2015) \u2014 SIR-style "
                     "compartmental models of hoax diffusion", "b"),
                    ("Reliability labels", "t"),
                    ("NewsGuard, Media Bias / Fact Check \u2014 outlet-level "
                     "reliability scores; used downstream of classification, "
                     "rarely as the response in a stochastic model", "b"),
                    ("Our contribution", "t"),
                    ("Joint event-level study: counting process for $N(t)$ "
                     "(Part A) plus a shared-likelihood comparison of Hawkes "
                     "and Markov dynamics for $\\pi_{pos}(t)$ (Part B)", "b"),
                ], title_size=12, body_size=10, formula_size=11,
                line_spacing=1.10)

# ---------------- Slide 3: Dataset Overview ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Dataset Overview",
             "CC-News articles cross-joined with NewsGuard reliability scores")
add_rich_bullets(s, Inches(0.45), Inches(1.15), Inches(6.4), Inches(5.6), [
    "Source: CC-News crawl, filtered to articles covering the 2017 Las Vegas shooting",
    "Date range: 2017-10-02 \u2192 2017-12-31 (91 days, zero-filled)",
    "2,702 articles  \u00b7  29.7 articles/day mean  \u00b7  peak 489 on day 0",
    "Reliability label constructed by joining CC-News with NewsGuard scores",
    ("Score thresholded at the median to produce a binary label $\\pm 1$", 1),
    ("63.5% reliable  /  36.5% unreliable across the corpus", 1),
    "Limitations of the dataset",
    ("Binary reliability label \u2014 the underlying NewsGuard score is continuous", 1),
    ("Daily granularity only \u2014 within-day burst structure is invisible", 1),
    ("Single event window \u2014 generality across event types is untested", 1),
    ("Coverage gaps on quiet days are zero-filled, not imputed", 1),
], size=15)
add_image_fit(s, PARTB / "fig1_observed_distribution.png",
              Inches(7.0), Inches(1.2), Inches(6.0), Inches(5.7))
add_plain_text(s, Inches(7.0), Inches(6.85), Inches(6.0), Inches(0.3),
               "Stacked daily counts by reliability category",
               size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# ---------------- Slide 4: Part A base models ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Part A  \u00b7  Modelling the Counting Process",
             "Three base models fit by maximum conditional Poisson "
             "log-likelihood")
add_rich_bullets(s, Inches(0.45), Inches(1.15), Inches(6.3), Inches(3.0), [
    "Standard Poisson  \u2014  $\\lambda(t) = \\lambda$  (1 param, null model)",
    "Inhomogeneous Poisson  \u2014  $\\lambda(t) = A\\,e^{-\\beta t} + c$  (3 params)",
    "Discrete-time Hawkes  \u2014  $\\lambda(t) = \\mu + \\alpha\\,R(t)$,  "
    "$R(t) = e^{-\\beta}\\,[R(t-1) + N(t-1)]$",
    "$\\beta$-profile sweep collapses to AR(1) limit as $\\beta \\to \\infty$",
    ("Optimal kernel:  $\\lambda(t) = 11.50 + 0.614\\,N(t-1)$", 1),
], size=14)

make_table(s, Inches(0.45), Inches(4.20), Inches(6.3), Inches(2.3),
           ["Model", "$k$", "$\\log L$", "AIC", "Pearson Var"],
           [
               ["Standard Poisson", "1", "$-3835.11$", "7672.22", "196.3"],
               ["Inhomog. Poisson", "3", "$-513.15$", "1032.29", "8.52"],
               ["Hawkes (AR-1)", "2", "$-2003.81$", "4011.62", "225.1"],
           ],
           highlight_row=2, header_size=12, body_size=11)
add_rich_paragraph(s, Inches(0.45), Inches(6.70), Inches(6.3), Inches(0.4),
                   "IHP wins decisively  \u2014  absorbs the day-0 exogenous "
                   "shock the Hawkes cannot.",
                   size=12, bold=True, color=NAVY)
add_image_fit(s, PARTA / "fig11_all_models_overlay.png",
              Inches(7.0), Inches(1.15), Inches(6.0), Inches(3.1))
add_image_fit(s, PARTA / "fig8_hawkes_beta_profile.png",
              Inches(7.0), Inches(4.30), Inches(6.0), Inches(2.85))

# ---------------- Slide 5: Part A hybrid winner ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Part A  \u00b7  Hybrid IHP + Hawkes Wins",
             "Combining exogenous trigger with endogenous self-excitation")
add_rich_paragraph(s, Inches(0.45), Inches(1.05), Inches(6.4), Inches(0.55),
                   "Hybrid (full):   "
                   "$\\lambda(t) = A\\,e^{-\\beta_{0} t} + c + \\alpha\\,R(t)$",
                   size=15, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(1.65), Inches(6.4), Inches(2.8), [
    "IHP term absorbs the day-0 burst (489 articles)",
    "Hawkes term carries the self-exciting tail",
    "Two timescales: fast exogenous decay (half-life $\\approx 1.75$ d) + "
    "slow endogenous persistence (half-life $\\approx 8.2$ d)",
    "LRT  IHP $\\to$ Hybrid-AR1:  $\\Delta\\ell = 124.4$,  $p \\approx 0$",
    "LRT  Hybrid-AR1 $\\to$ Hybrid-full:  $p = 3.7\\times 10^{-9}$",
], size=13)

make_table(s, Inches(0.45), Inches(4.55), Inches(6.4), Inches(2.35),
           ["Model", "$k$", "AIC", "BIC"],
           [
               ["Standard Poisson", "1", "7672.22", "7674.73"],
               ["Inhomog. Poisson", "3", "1032.29", "1039.82"],
               ["Hawkes (AR-1)", "2", "4011.62", "4016.65"],
               ["Hybrid IHP+Hawkes (AR-1)", "4", "785.59", "795.63"],
               ["Hybrid IHP+Hawkes (full)", "5", "752.80", "765.35"],
           ],
           highlight_row=5, header_size=11, body_size=10)

add_image_fit(s, PARTA / "fig13_hybrid_decomposition.png",
              Inches(7.0), Inches(1.10), Inches(6.0), Inches(3.0))
add_image_fit(s, PARTA / "fig14_hybrid_early_zoom.png",
              Inches(7.0), Inches(4.20), Inches(6.0), Inches(2.95))

# ---------------- Slide 6: Part B  setup + shared binomial LL ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Part B  \u00b7  Setup and the Shared Binomial Likelihood",
             "Conditioning on $N_{tot}(t)$ puts Hawkes and Markov on a common scale")
add_rich_bullets(s, Inches(0.45), Inches(1.05), Inches(6.15), Inches(6.0), [
    "Goal: model the reliable / unreliable split, not the total volume",
    ("$q_{t} = P(\\text{article on day } t \\text{ is reliable})$", 1),
    "Daily observations:  $N_{neg}(t),\\, N_{pos}(t)$;  "
    "$N_{tot}(t) = N_{neg}(t) + N_{pos}(t)$;  "
    "$\\pi_{pos}(t) = N_{pos}(t) / N_{tot}(t)$",
    "Why a shared likelihood is needed",
    ("Hawkes natively scores $(N_{neg}, N_{pos})$ via a Poisson log-likelihood", 1),
    ("Markov natively scores the transition $\\pi_{t-1} \\to \\pi_{t}$ via "
     "a multinomial", 1),
    ("Native AICs are on different scales \u2014 not directly comparable", 1),
    "Resolution: condition on the observed total $N_{tot}(t)$ and score the "
    "composition only",
    ("Justified by the Poisson\u2013binomial superposition identity (right)", 1),
    ("Semantically correct: Part B is about the *distribution*, not the volume", 1),
    "Constant-in-$q_t$ terms (the $\\log \\binom{N_{tot}}{N_{pos}}$ factor) "
    "cancel in every model comparison",
], size=13)

# Three stacked formula boxes on the right
RCOL_L = Inches(6.80); RCOL_W = Inches(6.30)
# Box 1 — Superposition identity
add_formula_box(s, RCOL_L, Inches(1.05), RCOL_W, Inches(1.70),
                "Poisson \u2013 binomial superposition",
                [
                    ("If  $X \\sim \\mathrm{Pois}(\\lambda_{1})$  and  "
                     "$Y \\sim \\mathrm{Pois}(\\lambda_{2})$  independently,", "t"),
                    ("$X \\mid (X + Y = n) \\;\\sim\\; "
                     "\\mathrm{Binomial}\\left(n,\\; "
                     "\\frac{\\lambda_{1}}{\\lambda_{1} + \\lambda_{2}}\\right)$", "f"),
                    ("\u21d2 applied to $(N_{neg}, N_{pos})$, this is "
                     "the binomial that scores the split", "sub"),
                ], title_size=12, body_size=11, formula_size=13)

# Box 2 — Per-day and total log-likelihood
add_formula_box(s, RCOL_L, Inches(2.90), RCOL_W, Inches(1.90),
                "Conditional binomial log-likelihood",
                [
                    ("Per day  $t$  with  $N_{tot}(t) > 0$:", "t"),
                    ("$\\ell_{t} \\;=\\; \\log \\binom{N_{tot}(t)}{N_{pos}(t)} "
                     "\\,+\\, N_{pos}(t) \\log q_{t} "
                     "\\,+\\, N_{neg}(t) \\log (1 - q_{t})$", "f"),
                    ("Total:  "
                     "$\\ell = \\sum_{t:\\, N_{tot}(t) > 0} \\ell_{t}$,  "
                     "effective sample size  $n_{eff} = 86$", "f"),
                ], title_size=12, body_size=11, formula_size=13)

# Box 3 — model-specific q_t
add_formula_box(s, RCOL_L, Inches(4.95), RCOL_W, Inches(2.15),
                "Model-specific reliable fraction  $q_{t}$",
                [
                    ("Hawkes (Poisson superposition):", "t"),
                    ("$q_{t} \\;=\\; \\dfrac{\\lambda_{pos}(t)}"
                     "{\\lambda_{pos}(t) + \\lambda_{neg}(t)}$", "f"),
                    ("Markov (one-step transition):", "t"),
                    ("$q_{t} \\;=\\; \\big(\\pi_{t-1}\\, P_{t}\\big)_{pos}"
                     " \\;=\\; \\pi_{neg}(t-1)\\, P_{neg \\to pos}"
                     " \\,+\\, \\pi_{pos}(t-1)\\, P_{pos \\to pos}$", "f"),
                ], title_size=12, body_size=11, formula_size=12)

# ---------------- Slide 7: Part B Model Specifications ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Part B  \u00b7  Model Specifications",
             "Four candidates on the shared binomial scale \u2014 "
             "full functional forms and derivations")

BOX_W = Inches(6.25); BOX_H = Inches(2.95)
L1 = Inches(0.40); L2 = Inches(6.80)
T1 = Inches(1.00); T2 = Inches(4.08)

# 1. Multivariate Hawkes
add_formula_box(s, L1, T1, BOX_W, BOX_H,
                "1.  Multivariate Hawkes  (AR-1 limit,  6 params)",
                [
                    ("$\\lambda_{k}(t) \\;=\\; \\mu_{k} \\;+\\; "
                     "\\sum_{j \\in \\{neg,\\, pos\\}} n_{k,j}\\, N_{j}(t-1),"
                     "\\quad k \\in \\{neg,\\, pos\\}$", "f"),
                    ("Parameters: baselines $\\mu_{neg},\\, \\mu_{pos}$  and "
                     "$2\\times 2$ branching matrix $n$", "b"),
                    ("Stationarity constraint: spectral radius "
                     "$\\rho(n) < 1$, enforced during optimisation", "b"),
                    ("Derivation: direct 2-D lift of Part A's AR(1) Hawkes; "
                     "$\\beta \\to \\infty$ collapse inherited", "b"),
                    ("Implied split:  "
                     "$q_{t} = \\lambda_{pos}(t) / "
                     "(\\lambda_{pos}(t) + \\lambda_{neg}(t))$", "b"),
                ], title_size=11, body_size=10, formula_size=12)

# 2. Homogeneous Markov
add_formula_box(s, L2, T1, BOX_W, BOX_H,
                "2z.  Homogeneous Markov  (null,  2 params)",
                [
                    ("$P \\;=\\; \\begin{pmatrix} P_{neg\\to neg} & "
                     "P_{neg \\to pos} \\\\ P_{pos \\to neg} & "
                     "P_{pos \\to pos} \\end{pmatrix},"
                     "\\quad \\text{rows sum to } 1$", "t"),
                    ("$q_{t} \\;=\\; (\\pi_{t-1}\\, P)_{pos} \\;=\\; "
                     "\\pi_{neg}(t-1)\\, P_{neg \\to pos} "
                     "\\,+\\, \\pi_{pos}(t-1)\\, P_{pos \\to pos}$", "f"),
                    ("Two free parameters (one per row after the simplex "
                     "constraint)", "b"),
                    ("Derivation: assume article-level first-order Markov; "
                     "aggregate to daily proportions", "b"),
                    ("Stationary $\\pi^{*}$ solves $\\pi^{*} P = \\pi^{*}"
                     " \\approx (0.39,\\, 0.61)$  \u2014  null benchmark", "b"),
                ], title_size=11, body_size=10, formula_size=11)

# 3. State-dependent Markov
add_formula_box(s, L1, T2, BOX_W, BOX_H,
                "2a.  State-dependent Markov  (regime-switching,  4 params)",
                [
                    ("$x_{t-1} \\;=\\; \\argmax_{k}\\, \\pi_{k}(t-1) "
                     "\\in \\{neg,\\, pos\\}$", "f"),
                    ("$q_{t} \\;=\\; \\big(\\pi_{t-1}\\, "
                     "P^{(x_{t-1})}\\big)_{pos},"
                     "\\quad P^{(neg)},\\, P^{(pos)} \\text{ separately fit}$",
                     "f"),
                    ("Two regime-specific transition matrices, keyed on "
                     "yesterday's majority class", "b"),
                    ("Derivation: treat the shock regime (neg-majority, "
                     "$\\le 12$ days) as a discrete latent state", "b"),
                    ("Threshold at $\\pi_{pos} = 0.5$ is model-imposed and "
                     "abrupt \u2014 smoothed by 2b below", "b"),
                ], title_size=11, body_size=10, formula_size=11)

# 4. Mean-field Markov
add_formula_box(s, L2, T2, BOX_W, BOX_H,
                "2b.  Mean-field Markov  (logistic links,  4 params)",
                [
                    ("$P_{neg \\to pos}(t) \\;=\\; "
                     "\\sigma\\!\\big(a_{0} + a_{1}\\, \\pi_{pos}(t-1)\\big)$",
                     "f"),
                    ("$P_{pos \\to neg}(t) \\;=\\; "
                     "\\sigma\\!\\big(b_{0} + b_{1}\\, \\pi_{pos}(t-1)\\big)$",
                     "f"),
                    ("Remaining entries from the row-sum constraint;  "
                     "$\\sigma(u) = 1/(1 + e^{-u})$", "b"),
                    ("Derivation: continuous analogue of 2a \u2014 regime "
                     "dependence becomes a smooth function of the empirical "
                     "proportion", "b"),
                    ("Nests the null at $a_{1} = b_{1} = 0$  \u21d2  clean "
                     "LRT, df $= 2$;  precursor to McKean\u2013Vlasov", "b"),
                ], title_size=11, body_size=10, formula_size=11)

# ---------------- Slide 8: Part B Multivariate Hawkes ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Part B  \u00b7  Multivariate Hawkes \u2014 Cross-Excitation",
             "Branching matrix between reliable and unreliable streams")
add_rich_bullets(s, Inches(0.45), Inches(1.15), Inches(6.4), Inches(2.8), [
    "Baselines:  $\\mu_{neg} = 3.12$,  $\\mu_{pos} = 8.38$",
    "Spectral radius  $\\rho(n) = 0.633$  $\\to$  stationary",
    "Within-category excitation:  $n_{nn} = 0.44$,  $n_{pp} = 0.37$",
    "Cross-excitation is non-trivial:  $n_{pn} = 0.32$,  $n_{np} = 0.16$",
    "Reliability is not a closed ecosystem \u2014 streams interact",
], size=14)
add_rich_paragraph(s, Inches(0.45), Inches(4.0), Inches(6.4), Inches(0.5),
                   "Justification: Hawkes is the natural multivariate "
                   "generalisation of Part A's winner;", size=12, color=SLATE)
add_rich_paragraph(s, Inches(0.45), Inches(4.35), Inches(6.4), Inches(0.5),
                   "branching matrix gives mechanistic, interpretable "
                   "cross-excitation estimates.", size=12, color=SLATE)
add_rich_paragraph(s, Inches(0.45), Inches(5.05), Inches(6.4), Inches(0.5),
                   "Interpretation", size=14, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(5.45), Inches(6.4), Inches(1.8), [
    "An unreliable article today triggers $\\approx 0.32$ reliable articles tomorrow",
    "A reliable article today triggers $\\approx 0.16$ unreliable articles tomorrow",
    "Mildly diagonal-dominant, sub-critical, and consistent with Part A",
], size=12)
add_image_fit(s, PARTB / "fig3_hawkes_branching_heatmap.png",
              Inches(7.05), Inches(1.10), Inches(6.0), Inches(3.0))
add_image_fit(s, PARTB / "fig4_hawkes_fitted_vs_observed.png",
              Inches(7.05), Inches(4.20), Inches(6.0), Inches(2.95))

# ---------------- Slide 9: Part B Markov variants ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Part B  \u00b7  Markov Variants \u2014 Regime Switching & "
                "Mean-Field",
             "Two parameterisations of state-dependent transition probabilities")
add_rich_paragraph(s, Inches(0.45), Inches(1.05), Inches(6.4), Inches(0.4),
                   "State-dependent (regime-switching) Markov",
                   size=14, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(1.45), Inches(6.4), Inches(2.0), [
    "Two transition matrices keyed on yesterday's majority state",
    "Pos-majority regime: both classes sticky "
    "($P_{nn} = 0.68$,  $P_{pp} = 0.76$)",
    "Neg-majority regime: dynamics push toward reliable "
    "($P_{neg \\to pos} = 0.73$)",
], size=12)
add_rich_paragraph(s, Inches(0.45), Inches(3.50), Inches(6.4), Inches(0.4),
                   "Mean-field Markov", size=14, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(3.90), Inches(6.4), Inches(2.9), [
    "Logistic links:  $P_{neg \\to pos} = \\sigma(a_{0} + a_{1}\\,\\pi_{pos})$,  "
    "$P_{pos \\to neg} = \\sigma(b_{0} + b_{1}\\,\\pi_{pos})$",
    "Fitted:  $a_{0} = +1.15$,  $a_{1} = -3.96$  $\\to$  strong mean reversion",
    "When reliable share $\\downarrow$, conversion to reliable $\\uparrow$ "
    "sharply ($0.76 \\to 0.06$)",
    "Smooth, identifiable analogue of the regime-switching matrix",
], size=12)
add_image_fit(s, PARTB / "fig5_markov_regime_matrices.png",
              Inches(7.05), Inches(1.10), Inches(6.0), Inches(3.0))
add_image_fit(s, PARTB / "fig6_meanfield_P_trajectory.png",
              Inches(7.05), Inches(4.20), Inches(6.0), Inches(2.95))

# ---------------- Slide 10: Part B comparison ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Part B  \u00b7  Model Comparison on Shared Binomial Scale",
             "Effective sample size $n_{eff} = 86$ days with at least "
             "one article")

make_table(s, Inches(0.45), Inches(1.15), Inches(6.4), Inches(2.4),
           ["Model", "$k$", "$\\log L$", "AIC", "BIC"],
           [
               ["Multivariate Hawkes", "6", "$-156.93$", "325.87", "340.60"],
               ["Homogeneous Markov", "2", "$-162.75$", "329.50", "334.41"],
               ["State-dependent Markov", "4", "$-157.33$", "322.65", "332.47"],
               ["Mean-field Markov", "4", "$-157.56$", "323.13", "332.94"],
           ],
           highlight_row=3, header_size=12, body_size=11)

add_rich_paragraph(s, Inches(0.45), Inches(3.7), Inches(6.4), Inches(0.4),
                   "Likelihood ratio tests vs homogeneous baseline",
                   size=13, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(4.05), Inches(6.4), Inches(1.4), [
    "State-dep Markov:  $\\mathrm{LR} = 10.85$,  df $= 2$,  "
    "$p = 4.4\\times 10^{-3}$",
    "Mean-field Markov:  $\\mathrm{LR} = 10.37$,  df $= 2$,  "
    "$p = 5.6\\times 10^{-3}$",
], size=12)
add_rich_paragraph(s, Inches(0.45), Inches(5.45), Inches(6.4), Inches(0.4),
                   "Why Markov beats Hawkes here",
                   size=14, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(5.85), Inches(6.4), Inches(1.4), [
    "Hawkes spends parameters on total counts, not the split",
    "$\\sim 86\\%$ of days are pos-majority \u2014 a single $P$ captures most dynamics",
], size=12)

add_image_fit(s, PARTB / "fig7_fitted_proportion_comparison.png",
              Inches(7.05), Inches(1.10), Inches(6.0), Inches(3.0))
add_image_fit(s, PARTB / "fig9_model_comparison.png",
              Inches(7.05), Inches(4.20), Inches(6.0), Inches(2.95))

# ---------------- Slide 11: Part B interpretation ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Part B  \u00b7  Interpretation \u2014 A Mean-Reverting "
                "Reliability Ecology",
             "What the winning models say about the news ecosystem")
add_rich_bullets(s, Inches(0.45), Inches(1.15), Inches(12.4), Inches(5.5), [
    "The system has a stable equilibrium near $\\pi_{pos} \\approx 0.63$ "
    "\u2014 the empirical reliable share",
    "When the reliable share drops below $0.5$ (only at the initial shock), "
    "strong restoring forces fire",
    ("State-dependent Markov:  $P_{neg \\to pos}$ jumps to $0.73$ in the "
     "neg-majority regime", 1),
    ("Mean-field Markov:  $P_{neg \\to pos}$ is a steeply decreasing "
     "function of $\\pi_{pos}$", 1),
    "Hawkes branching matrix encodes the same story mechanistically",
    ("Reliable articles preferentially trigger more reliable articles  "
     "($n_{pp} > n_{np}$)", 1),
    ("Unreliable articles trigger reliable and unreliable articles nearly "
     "symmetrically", 1),
    "Choice of model depends on the question being asked",
    ("Distribution dynamics  $\\to$  Markov variants are most efficient", 1),
    ("Mechanistic cross-excitation  $\\to$  Multivariate Hawkes is most "
     "informative", 1),
    "Synthetic recovery passed (max rel. error $20\\%$ at $T = 500$) \u2014 "
    "estimates are trustworthy",
], size=15)

# ---------------- Slide 12: Conclusions + next steps ----------------
s = prs.slides.add_slide(BLANK)
slide_header(s, "Conclusions and Next Steps")
add_rich_paragraph(s, Inches(0.45), Inches(1.05), Inches(6.2), Inches(0.45),
                   "Summary of findings",
                   size=16, bold=True, color=NAVY)
add_rich_bullets(s, Inches(0.45), Inches(1.50), Inches(6.2), Inches(5.5), [
    "Counting process is best modelled by a hybrid IHP + Hawkes",
    ("Exogenous shock + endogenous amplification at distinct timescales", 1),
    "Reliability split is best modelled by a state-dependent Markov chain",
    ("Mean-field Markov is statistically indistinguishable and smoother", 1),
    "Cross-excitation between reliable and unreliable streams is non-trivial",
    "Reliability ecology is mean-reverting around the empirical equilibrium",
], size=14)
add_rich_paragraph(s, Inches(7.0), Inches(1.05), Inches(6.0), Inches(0.45),
                   "Next steps", size=16, bold=True, color=NAVY)
add_rich_bullets(s, Inches(7.0), Inches(1.50), Inches(6.0), Inches(5.5), [
    "McKean\u2013Vlasov formulation for continuous-time, "
    "distribution-dependent transitions",
    ("Replace discrete $P_{t}$ with an SDE driven by the empirical "
     "$\\pi_{pos}$ measure", 1),
    "Use raw NewsGuard scores (continuous) instead of the binary label",
    ("Recovers within-class heterogeneity lost to the median threshold", 1),
    ("Enables regression / functional models on a real-valued reliability axis", 1),
    "Replicate across additional events to test generality of the "
    "hybrid + Markov story",
    "Extend hybrid model with explicit weekly seasonality "
    "($\\Delta \\mathrm{AIC} \\approx -148$ in exploration)",
], size=14)

prs.save(OUT)
print(f"Wrote {OUT}")
